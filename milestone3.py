import streamlit as st
import nltk
import re  # for email validation
from db1 import register_user, login_user, save_profile, get_all_users, log_request, get_logs

# Summarizer (using sumy for simple extractive summaries)
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer

# Readability analysis
import textstat
import matplotlib.pyplot as plt

# Transformers summarizers + ROUGE
from transformers import pipeline
import evaluate

# File parsing libraries
import docx
from pptx import Presentation
from PyPDF2 import PdfReader

# HuggingFace dataset
from datasets import load_dataset

# ----------------- Email Validator -----------------
def is_valid_email(email):
    return re.match(r"[^@]+@[^@]+\.[^@]+", email) is not None

# ----------------- Summarization Function (Sumy) -----------------
def summarize_text(text, sentences_count=3):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer()
    summary = summarizer(parser.document, sentences_count)
    return " ".join(str(sentence) for sentence in summary)

# ----------------- HuggingFace Summarizers -----------------
@st.cache_resource(show_spinner=False)
def get_summarizer(model_key: str):
    model_map = {
        "DistilBART (CNN/DM, small)": "sshleifer/distilbart-cnn-12-6",
        "T5-Small": "t5-small",
        "BART-Base": "facebook/bart-base",
    }
    return pipeline("summarization", model=model_map[model_key])

LENGTH_PRESETS = {
    "Short": {"min_length": 20, "max_length": 40},
    "Medium": {"min_length": 60, "max_length": 100},
    "Long": {"min_length": 120, "max_length": 180},
}

@st.cache_resource(show_spinner=False)
def get_rouge():
    return evaluate.load("rouge")

# ----------------- Load Dataset -----------------
@st.cache_resource(show_spinner=False)
def get_dataset():
    return load_dataset("cnn_dailymail", "3.0.0")

# ----------------- Extract Text from Different File Types -----------------
def extract_text_from_file(uploaded_file):
    if uploaded_file.name.endswith(".txt"):
        return uploaded_file.read().decode("utf-8")
    elif uploaded_file.name.endswith(".pdf"):
        reader = PdfReader(uploaded_file)
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
    elif uploaded_file.name.endswith(".docx"):
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    elif uploaded_file.name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    return None

# ----------------- Custom CSS -----------------
st.markdown(
    """
    <style>
    .card {
        background: #ffffff;
        padding: 20px;
        border-radius: 12px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        margin: 10px;
    }
    .stButton>button {
        background-color: #00796b;
        color: white;
        border-radius: 8px;
        padding: 8px 20px;
        border: none;
    }
    .stButton>button:hover {
        background-color: #004d40;
        color: white;
    }
    h3 { color: #004d40; }
    </style>
    """,
    unsafe_allow_html=True,
)

# ----------------- Page Title -----------------
st.title("Secure Authentication, Summarization, Paraphrasing & Dataset Management")

col1, col2 = st.columns(2)

# ----------------- User Authentication -----------------
with col1:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("🔑 User Authentication")

    email = st.text_input("Email", placeholder="user@example.com")
    password = st.text_input("Password", type="password", placeholder="Enter password")

    col_a, col_b = st.columns([1, 1])
    with col_a:
        if st.button("Sign In"):
            if login_user(email, password):
                st.session_state["logged_in"] = email
                st.success("✅ Login successful!")
            else:
                st.error("❌ Invalid credentials")

    with col_b:
        if st.button("Create Account"):
            if not is_valid_email(email):
                st.error("⚠️ Please enter a valid email address")
            elif not password:
                st.error("⚠️ Password cannot be empty")
            else:
                try:
                    register_user(email, password)
                    st.success("✅ Account created successfully!")
                except:
                    st.error("⚠️ User already exists")

    st.markdown("<a href='#'>Forgot password?</a>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Profile Management -----------------
with col2:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("👤 Profile Management")

    if "logged_in" in st.session_state:
        name = st.text_input("Name", placeholder="Your name")
        age_group = st.selectbox("Age Group", ["18-25", "26-35", "36-50", "50+"])
        st.write("Language Preference:")
        lang = st.radio("Select Language", ["English", "हिंदी"], horizontal=True)

        if st.button("Save Profile"):
            save_profile(st.session_state["logged_in"], name, age_group, lang)
            st.success("✅ Profile saved successfully!")
    else:
        st.warning("⚠️ Please login first to manage profile")

    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Show Database Table -----------------
if st.checkbox("Show Database Table"):
    rows = get_all_users()
    st.subheader("Users Table")
    st.table(rows)

# ----------------- Dashboard: Document Upload & Summarization + Readability -----------------
if "logged_in" in st.session_state:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("📊 Dashboard: Summarization & Readability Analysis")

    uploaded_file = st.file_uploader("📂 Upload a document", type=["txt", "pdf", "docx", "pptx"])
    text_input = st.text_area("✍️ Or paste text here:")

    if st.button("🔍 Analyze Text"):
        if uploaded_file is not None:
            text = extract_text_from_file(uploaded_file)
        elif text_input.strip() != "":
            text = text_input
        else:
            text = None
            st.warning("⚠️ Please upload a file or paste some text")

        if text:
            # Show original text
            st.subheader("📑 Original Content")
            st.write(text)

            # Summarization
            st.subheader("📝 Summarized Content (Sumy)")
            summary = summarize_text(text)
            st.success(summary)

            # ✅ Readability Analysis with Chart
            st.subheader("📈 Readability Scores")
            flesch = textstat.flesch_kincaid_grade(text)
            fog = textstat.gunning_fog(text)
            smog = textstat.smog_index(text)

            scores = [flesch, fog, smog]
            labels = ["Flesch-Kincaid", "Gunning Fog", "SMOG"]

            fig, ax = plt.subplots()
            bars = ax.bar(labels, scores, color=["#4CAF50", "#FF9800", "#2196F3"])
            ax.set_ylabel("Score")
            ax.set_title("Readability Analysis")

            # Annotate bars
            for bar, score in zip(bars, scores):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height(),
                        f"{score:.2f}", ha="center", va="bottom", fontsize=10)

            st.pyplot(fig)
    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Pretrained Summarization + ROUGE -----------------
if "logged_in" in st.session_state:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("🧠 Pretrained Summarization + ROUGE")

    model_choice = st.selectbox("Choose model", ["DistilBART (CNN/DM, small)", "T5-Small", "BART-Base"])
    length_choice = st.radio("Summary Length", ["Short", "Medium", "Long"], horizontal=True)

    src_col, ref_col = st.columns(2)
    with src_col:
        doc_text = st.text_area("✍️ Paste text to summarize", height=220)
    with ref_col:
        ref_summary = st.text_area("🎯 (Optional) Reference summary", height=220)

    if st.button("⚡ Generate Pretrained Summary"):
        if doc_text.strip():
            with st.spinner("Generating summary..."):
                summarizer = get_summarizer(model_choice)
                params = LENGTH_PRESETS[length_choice]
                input_text = f"summarize: {doc_text}" if "T5" in model_choice else doc_text
                output = summarizer(
                    input_text,
                    min_length=params["min_length"],
                    max_length=params["max_length"],
                    truncation=True,
                    do_sample=False,
                )
                summary_text = output[0]["summary_text"]

            st.subheader("📝 Pretrained Summary")
            st.success(summary_text)

            # ✅ Log
            log_request(
                st.session_state.get("logged_in", "guest"),
                doc_text,
                "summarization",
                model_choice,
                summary_text,
            )

            if ref_summary.strip():
                rouge = get_rouge()
                scores = rouge.compute(predictions=[summary_text], references=[ref_summary])
                st.subheader("📊 ROUGE Scores")
                st.json(scores)
        else:
            st.warning("⚠️ Please enter text to summarize")
    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Dataset Management -----------------
if "logged_in" in st.session_state:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("📚 Dataset Summarization (CNN/DailyMail)")

    dataset = get_dataset()
    split_choice = st.selectbox("Choose dataset split", ["train", "validation", "test"])
    sample_id = st.number_input("Sample ID", min_value=0, max_value=1000, value=0)

    item = dataset[split_choice][sample_id]
    original_text = item["article"]
    reference_text = item["highlights"]

    st.write("### 📑 Original Text")
    st.info(original_text[:1000] + "..." if len(original_text) > 1000 else original_text)

    st.write("### 🎯 Reference Summary")
    st.success(reference_text)

    if st.button("⚡ Generate Model Summary (Dataset)"):
        summarizer = get_summarizer("DistilBART (CNN/DM, small)")
        summary = summarizer(original_text, max_length=120, min_length=30, do_sample=False)[0]["summary_text"]

        st.write("### 📝 Model Summary")
        st.success(summary)

        rouge = get_rouge()
        scores = rouge.compute(predictions=[summary], references=[reference_text])
        st.write("### 📊 ROUGE Scores")
        st.json(scores)

        log_request(
            st.session_state.get("logged_in", "guest"),
            original_text,
            "dataset_summarization",
            "DistilBART",
            summary,
        )
    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Paraphrasing -----------------
if "logged_in" in st.session_state:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("📝 Paraphrasing with FLAN-T5 / BART")

    model_choice_para = st.selectbox("Choose paraphrasing model", ["FLAN-T5 (small)", "BART (base)"])
    complexity_choice = st.selectbox("Select complexity level", ["Simple", "Moderate", "Advanced"])
    para_text = st.text_area("✍️ Enter text to paraphrase", height=200)

    if st.button("🔄 Generate Paraphrase"):
        if para_text.strip():
            with st.spinner("Generating paraphrase..."):
                model_map_para = {"FLAN-T5 (small)": "google/flan-t5-small", "BART (base)": "facebook/bart-base"}
                paraphraser = pipeline("text2text-generation", model=model_map_para[model_choice_para])

                prompt = (
                    f"Rephrase the following text in a {complexity_choice.lower()} way, "
                    f"using different wording but keeping the meaning same:\n\n{para_text}"
                )
                result = paraphraser(
                    prompt, max_length=220, do_sample=True, top_k=50, top_p=0.95, temperature=0.9
                )

                para_output = result[0].get("summary_text") or result[0].get("generated_text") or str(result[0])

            col1, col2 = st.columns(2)
            with col1:
                st.subheader("📑 Original")
                st.info(para_text)
            with col2:
                st.subheader("🔄 Paraphrased")
                st.success(para_output)

            log_request(
                st.session_state.get("logged_in", "guest"),
                para_text,
                "paraphrasing",
                model_choice_para,
                para_output,
            )
        else:
            st.warning("⚠️ Please enter text to paraphrase")
    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Logs Viewer -----------------
if "logged_in" in st.session_state:
    if st.checkbox("📜 Show My Logs"):
        logs = get_logs(st.session_state["logged_in"])
        st.subheader("User Logs")
        st.table(logs)

# ----------------- ByT5 Fine-Tuning & Inference -----------------
if "logged_in" in st.session_state:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("🤖 Fine-Tune ByT5 on CNN/DailyMail")

    st.info(
        "This section allows you to fine-tune the ByT5-small model on a subset of the CNN/DailyMail dataset. "
        "Training may take several minutes depending on your machine/GPU."
    )

    num_train = st.number_input("Number of training samples", min_value=100, max_value=5000, value=2000)
    num_val   = st.number_input("Number of validation samples", min_value=50, max_value=1000, value=500)
    num_epochs = st.number_input("Number of epochs", min_value=1, max_value=10, value=3)

    if st.button("⚡ Start ByT5 Training"):
        import torch
        from datasets import load_dataset
        from transformers import (
            AutoTokenizer, AutoModelForSeq2SeqLM,
            DataCollatorForSeq2Seq, Seq2SeqTrainer, Seq2SeqTrainingArguments
        )

        with st.spinner("Loading dataset..."):
            dataset = load_dataset("cnn_dailymail", "3.0.0")
            train_data = dataset["train"].shuffle(seed=42).select(range(num_train))
            val_data   = dataset["validation"].shuffle(seed=42).select(range(num_val))

        with st.spinner("Loading model and tokenizer..."):
            model_name = "google/byt5-small"
            tokenizer = AutoTokenizer.from_pretrained(model_name)
            model = AutoModelForSeq2SeqLM.from_pretrained(model_name)

        st.info("Preprocessing dataset...")

        max_input_len = 512
        max_output_len = 128

        def preprocess(batch):
            inputs = tokenizer(batch["article"], truncation=True, padding="max_length", max_length=max_input_len)
            with tokenizer.as_target_tokenizer():
                labels = tokenizer(batch["highlights"], truncation=True, padding="max_length", max_length=max_output_len)
            inputs["labels"] = labels["input_ids"]
            return inputs

        train_enc = train_data.map(preprocess, batched=True, remove_columns=["article","highlights","id"])
        val_enc   = val_data.map(preprocess, batched=True, remove_columns=["article","highlights","id"])

        st.info("Setting up training arguments...")
        training_args = Seq2SeqTrainingArguments(
            output_dir="./byt5_cnn",
            per_device_train_batch_size=4,
            per_device_eval_batch_size=4,
            num_train_epochs=num_epochs,
            learning_rate=5e-5,
            weight_decay=0.01,
            save_total_limit=2,
            predict_with_generate=True,
            fp16=torch.cuda.is_available(),
            logging_dir="./logs",
            logging_steps=50
        )

        data_collator = DataCollatorForSeq2Seq(tokenizer, model=model)

        st.info("Initializing trainer...")
        trainer = Seq2SeqTrainer(
            model=model,
            args=training_args,
            train_dataset=train_enc,
            eval_dataset=val_enc,
            tokenizer=tokenizer,
            data_collator=data_collator,
        )

        st.info("Training started! This may take a while... 🚀")
        trainer.train()
        st.success("✅ Training completed!")

        st.info("Evaluating model...")
        metrics = trainer.evaluate()
        st.json(metrics)

        st.info("Saving fine-tuned model...")
        trainer.save_model("./byt5_cnn")
        tokenizer.save_pretrained("./byt5_cnn")
        st.success("✅ Model saved at ./byt5_cnn")

        # ================================
        # Inference / Generate summaries
        # ================================
        st.subheader("✍️ Generate Summary")
        input_text = st.text_area("Enter article text to summarize:")

        if st.button("Generate Summary"):
            if input_text.strip() != "":
                st.info("Generating summary... 🚀")
                # Load model & tokenizer from saved folder
                model = AutoModelForSeq2SeqLM.from_pretrained("./byt5_cnn")
                tokenizer = AutoTokenizer.from_pretrained("./byt5_cnn")

                inputs = tokenizer(input_text, return_tensors="pt", truncation=True, padding="max_length", max_length=512)
                outputs = model.generate(**inputs, max_new_tokens=128)  # max_new_tokens used only here
                summary = tokenizer.decode(outputs[0], skip_special_tokens=True)

                st.success("✅ Summary generated!")
                st.write(summary)
            else:
                st.warning("Please enter some text to summarize.")

    st.markdown("</div>", unsafe_allow_html=True)
