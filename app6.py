
import streamlit as st
import nltk
import re  # for email validation
from db import register_user, login_user, save_profile, get_all_users

# Summarizer (using sumy for simple extractive summaries)
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer

# Readability analysis
import textstat
import matplotlib.pyplot as plt

# Transformers summarizers + ROUGE
from transformers import pipeline
import evaluate

# File parsing libraries
import docx
from pptx import Presentation
from PyPDF2 import PdfReader

# ----------------- Email Validator -----------------
def is_valid_email(email):
    return re.match(r"[^@]+@[^@]+\.[^@]+", email) is not None

# ----------------- Summarization Function (Sumy) -----------------
def summarize_text(text, sentences_count=3):
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LsaSummarizer()
    summary = summarizer(parser.document, sentences_count)
    return " ".join(str(sentence) for sentence in summary)

# ----------------- HuggingFace Summarizers (Week 4, lightweight models) -----------------
@st.cache_resource(show_spinner=False)
def get_summarizer(model_key: str):
    model_map = {
        "DistilBART (CNN/DM, small)": "sshleifer/distilbart-cnn-12-6",  # ~300MB
        "T5-Small": "t5-small",                                         # ~250MB
        "BART-Base": "facebook/bart-base",                              # ~500MB
    }
    return pipeline("summarization", model=model_map[model_key])

LENGTH_PRESETS = {
    "Short":  {"min_length": 20,  "max_length": 40},
    "Medium": {"min_length": 60,  "max_length": 100},
    "Long":   {"min_length": 120, "max_length": 180},
}

@st.cache_resource(show_spinner=False)
def get_rouge():
    return evaluate.load("rouge")

# ----------------- Extract Text from Different File Types -----------------
def extract_text_from_file(uploaded_file):
    if uploaded_file.name.endswith(".txt"):
        return uploaded_file.read().decode("utf-8")

    elif uploaded_file.name.endswith(".pdf"):
        reader = PdfReader(uploaded_file)
        return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

    elif uploaded_file.name.endswith(".docx"):
        doc = docx.Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    elif uploaded_file.name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    return None

# ----------------- Custom CSS -----------------
st.markdown("""
    <style>
    .card {
        background: #ffffff;
        padding: 20px;
        border-radius: 12px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        margin: 10px;
    }
    .stButton>button {
        background-color: #00796b;
        color: white;
        border-radius: 8px;
        padding: 8px 20px;
        border: none;
    }
    .stButton>button:hover {
        background-color: #004d40;
        color: white;
    }
    h3 {
        color: #004d40;
    }
    </style>
""", unsafe_allow_html=True)

# ----------------- Page Title -----------------
st.title("Milestone 1: User Authentication, Summarization & Readability Dashboard")

col1, col2 = st.columns(2)

# ----------------- User Authentication -----------------
with col1:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("🔑 User Authentication")

    email = st.text_input("Email", placeholder="user@example.com")
    password = st.text_input("Password", type="password", placeholder="Enter password")

    col_a, col_b = st.columns([1,1])
    with col_a:
        if st.button("Sign In"):
            if login_user(email, password):
                st.session_state['logged_in'] = email
                st.success("✅ Login successful!")
            else:
                st.error("❌ Invalid credentials")

    with col_b:
        if st.button("Create Account"):
            if not is_valid_email(email):
                st.error("⚠️ Please enter a valid email address")
            elif not password:
                st.error("⚠️ Password cannot be empty")
            else:
                try:
                    register_user(email, password)
                    st.success("✅ Account created successfully!")
                except:
                    st.error("⚠️ User already exists")

    st.markdown("<a href='#'>Forgot password?</a>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Profile Management -----------------
with col2:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("👤 Profile Management")

    if "logged_in" in st.session_state:
        name = st.text_input("Name", placeholder="Your name")
        age_group = st.selectbox("Age Group", ["18-25", "26-35", "36-50", "50+"])
        st.write("Language Preference:")
        lang = st.radio("Select Language", ["English", "हिंदी"], horizontal=True)

        if st.button("Save Profile"):
            save_profile(st.session_state['logged_in'], name, age_group, lang)
            st.success("✅ Profile saved successfully!")
    else:
        st.warning("⚠️ Please login first to manage profile")

    st.markdown("</div>", unsafe_allow_html=True)

# ----------------- Show Database Table -----------------
if st.checkbox("Show Database Table"):
    rows = get_all_users()
    st.subheader("Users Table")
    st.table(rows)

# ----------------- Dashboard: Document Upload & Summarization + Readability -----------------
if "logged_in" in st.session_state:
    st.markdown("<div class='card'>", unsafe_allow_html=True)
    st.subheader("📊 Dashboard: Summarization & Readability Analysis")

    uploaded_file = st.file_uploader("📂 Upload a document", type=["txt", "pdf", "docx", "pptx"])
    text_input = st.text_area("✍️ Or paste text here:")

    if st.button("🔍 Analyze Text"):
        if uploaded_file is not None:
            text = extract_text_from_file(uploaded_file)
        elif text_input.strip() != "":
            text = text_input
        else:
            text = None
            st.warning("⚠️ Please upload a file or paste some text")

        if text:
            # Show original text
            st.subheader("📑 Original Content")
            st.write(text)

            # Summarization
            st.subheader("📝 Summarized Content (Sumy)")
            summary = summarize_text(text)
            st.success(summary)

            # Readability Analysis
            st.subheader("📈 Readability Scores")
            flesch = textstat.flesch_kincaid_grade(text)
            fog = textstat.gunning_fog(text)
            smog = textstat.smog_index(text)

            col1, col2, col3 = st.columns(3)
            col1.metric("Flesch-Kincaid", f"{flesch:.2f}")
            col2.metric("Gunning Fog", f"{fog:.2f}")
            col3.metric("SMOG Index", f"{smog:.2f}")

            # Bar Chart
            scores = [flesch, fog, smog]
            labels = ["Beginner", "Intermediate", "Advanced"]
            normalized = [(s / max(scores)) * 100 for s in scores]

            fig, ax = plt.subplots()
            bars = ax.bar(labels, scores, color=["green", "orange", "red"])
            ax.set_ylabel("Score")
            ax.set_title("Readability Levels")

            for bar, score, norm in zip(bars, scores, normalized):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() - 1,
                        f"{score:.2f} ({norm:.1f}%)",
                        ha='center', color='white', fontsize=10, fontweight='bold')

            st.pyplot(fig)

    st.markdown("</div>", unsafe_allow_html=True)


